import streamlit as st
import google.generativeai as genai
from tavily import TavilyClient
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
import os
import json
import requests
import urllib.parse
from io import BytesIO
import random
import time

# --- 1. Configurare Pagină ---
st.set_page_config(page_title="Marketing AI (Hybrid)", page_icon="🛡️", layout="wide")

# --- 2. Secrete ---
try:
    GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
    TAVILY_API_KEY = st.secrets["TAVILY_API_KEY"]
    # Încearcă să ia cheia HF, dar nu oprește scriptul dacă lipsește (folosește fallback)
    HF_API_KEY = st.secrets.get("HF_API_KEY", "") 
except:
    st.error("⚠️ Lipsesc cheile API de bază (Google/Tavily). Verifică secrets.")
    st.stop()

genai.configure(api_key=GOOGLE_API_KEY)
tavily_client = TavilyClient(api_key=TAVILY_API_KEY)

# --- 3. Funcții Backend ---

@st.cache_data(ttl=3600)
def get_available_models():
    try:
        models = [m.name for m in genai.list_models() if 'generateContent' in m.supported_generation_methods and 'gemini' in m.name]
        return sorted(models, reverse=True)
    except:
        return ["models/gemini-1.5-flash"]

def upload_to_gemini(uploaded_file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
            tmp.write(uploaded_file.getvalue())
            tmp_path = tmp.name
        ref = genai.upload_file(tmp_path, mime_type="application/pdf")
        os.remove(tmp_path)
        return ref
    except: return None

def search_internet(query):
    try:
        res = tavily_client.search(query=query, search_depth="advanced", max_results=3)
        return "\n".join([f"- {r['content']}" for r in res.get('results', [])])
    except: return "Fără date."

# --- SISTEM HYBRID DE GENERARE IMAGINI ---

def try_pollinations(prompt_text):
    """Metoda de rezervă (Backup): Pollinations"""
    try:
        # Folosim un seed random și așteptăm puțin pentru a evita Rate Limit
        time.sleep(random.uniform(2, 4)) 
        
        short_prompt = prompt_text[:150]
        encoded_prompt = urllib.parse.quote(short_prompt)
        random_seed = random.randint(1, 100000)
        
        # Folosim modelul Turbo (mai rapid)
        url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&model=turbo&nologo=true&seed={random_seed}"
        
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers, timeout=20)
        
        if response.status_code == 200 and "image" in response.headers.get("Content-Type", ""):
            return BytesIO(response.content)
    except:
        pass
    return None

def try_huggingface(prompt_text):
    """Metoda Principală: Hugging Face (Stable Diffusion v1.5 - Rapid)"""
    if not HF_API_KEY:
        return None # Nu avem cheie, trecem la backup

    # Folosim v1-5 pentru că e "ușor" și răspunde rapid pe Free Tier
    API_URL = "https://api-inference.huggingface.co/models/runwayml/stable-diffusion-v1-5"
    headers = {"Authorization": f"Bearer {HF_API_KEY}"}
    payload = {"inputs": prompt_text}

    try:
        response = requests.post(API_URL, headers=headers, json=payload, timeout=15)
        if response.status_code == 200:
            return BytesIO(response.content)
    except:
        pass
    return None

def generate_image_robust(prompt_text):
    """Încearcă HF -> Dacă pică -> Încearcă Pollinations"""
    
    # 1. Încercăm Hugging Face
    print("Încerc Hugging Face...")
    img = try_huggingface(prompt_text)
    if img: return img
    
    # 2. Dacă nu a mers, încercăm Pollinations
    print("HF a eșuat. Încerc Pollinations...")
    img = try_pollinations(prompt_text)
    if img: return img
    
    return None

# --- 4. Generare PPT ---

def create_presentation_robust(slides_json, progress_callback=None):
    prs = Presentation()
    try:
        data = json.loads(slides_json)
    except:
        return None

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data.get("presentation_title", "Marketing")
    slide.placeholders[1].text = "Generat cu AI Hybrid"
    blank_layout = prs.slide_layouts[6] 
    total = len(data.get("slides", []))

    for i, slide_data in enumerate(data.get("slides", [])):
        if progress_callback:
            progress_callback(int((i/total)*90), f"Generez Slide {i+1}/{total} (Caut cea mai bună sursă foto)...")

        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes
        
        # Text
        tb = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        tb.text_frame.paragraphs[0].text = slide_data.get("title", "Slide")
        tb.text_frame.paragraphs[0].font.size = Pt(32)
        tb.text_frame.paragraphs[0].font.bold = True
        tb.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        tb_body = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        tf = tb_body.text_frame
        tf.word_wrap = True
        for point in slide_data.get("points", []):
            p = tf.add_paragraph()
            p.text = "• " + point
            p.font.size = Pt(18)
            p.space_after = Pt(12)

        # IMAGINE ROBUSTĂ
        image_prompt = slide_data.get("image_prompt", "")
        image_bytes = None
        
        if image_prompt:
            image_bytes = generate_image_robust(image_prompt)
        
        if image_bytes:
            try:
                shapes.add_picture(image_bytes, Inches(5.5), Inches(1.8), Inches(4.2), Inches(4.2))
            except:
                image_bytes = None
        
        if not image_bytes:
            # Fallback Vizual
            shape = shapes.add_shape(1, Inches(5.5), Inches(1.8), Inches(4.2), Inches(4.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(230, 230, 230)
            shape.text_frame.text = "Imagine Indisponibilă"
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        return tmp.name

# --- 5. UI ---

st.title("🛡️ Asistent Marketing (Sistem Robust)")

with st.sidebar:
    st.header("Setări")
    model_name = st.selectbox("Model", get_available_models(), format_func=lambda x: x.replace("models/", "").upper())
    uploaded_file = st.file_uploader("Catalog PDF", type=['pdf'])
    if st.button("Reset"):
        st.session_state.clear()
        st.rerun()

if "messages" not in st.session_state: st.session_state.messages = []

if uploaded_file and "gemini_file" not in st.session_state:
    with st.spinner("Procesez PDF..."):
        ref = upload_to_gemini(uploaded_file)
        if ref:
            st.session_state.gemini_file = ref
            st.success("PDF OK!")

for msg in st.session_state.messages:
    with st.chat_message(msg["role"]): st.markdown(msg["content"])

if prompt := st.chat_input("Ex: Idei produse"):
    if "gemini_file" not in st.session_state: st.error("Încarcă PDF.")
    else:
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"): st.markdown(prompt)

        with st.chat_message("assistant"):
            with st.spinner("Analizez..."):
                web = search_internet(prompt)
                try:
                    model = genai.GenerativeModel(model_name)
                    resp = model.generate_content([f"Context PDF + Net: {web}. Întrebare: {prompt}", st.session_state.gemini_file])
                    st.markdown(resp.text)
                    st.session_state.messages.append({"role": "assistant", "content": resp.text})
                    st.session_state.last_analysis = resp.text
                except Exception as e:
                    st.error(f"Eroare: {e}")

if "last_analysis" in st.session_state:
    st.divider()
    if st.button("✨ Generează Prezentare (Auto-Fix)"):
        
        progress_bar = st.progress(0, text="Structura...")
        
        try:
            try:
                json_model = genai.GenerativeModel(model_name, generation_config={"response_mime_type": "application/json"})
            except:
                json_model = genai.GenerativeModel(model_name)

            prompt_slides = f"""
            Pe baza analizei: {st.session_state.last_analysis}
            Generează JSON pentru 4-5 slide-uri.
            Include 'image_prompt' (ENGLEZĂ, max 20 cuvinte).
            
            FORMAT:
            {{
                "presentation_title": "Titlu",
                "slides": [
                    {{ "title": "Titlu", "points": ["Punct 1"], "image_prompt": "Red pen on desk" }}
                ]
            }}
            """
            
            resp = json_model.generate_content(prompt_slides)
            json_text = resp.text.replace("```json", "").replace("```", "").strip()
            
            pptx_path = create_presentation_robust(json_text, progress_bar.progress)
            
            progress_bar.progress(100, text="Gata!")
            
            if pptx_path:
                with open(pptx_path, "rb") as f:
                    st.download_button("📥 Descarcă PPTX", f, "Prezentare_Robust.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        except Exception as e:
            st.error(f"Eroare: {e}")
